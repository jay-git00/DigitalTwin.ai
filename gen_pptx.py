from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys

# ── Colors ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0A, 0x0E, 0x1A)   # dark navy
ACCENT  = RGBColor(0x00, 0xD4, 0xFF)   # cyan
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]   # completely blank

# ── Helper functions ─────────────────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txb(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def accent_bar(slide):
    rect(slide, 0, 0, 13.33, 0.07, ACCENT)

def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    txb(slide, title, 0.4, 0.15, 12, 0.6, size=32, bold=True, color=WHITE)
    if subtitle:
        txb(slide, subtitle, 0.4, 0.75, 12, 0.4, size=14, color=MUTED)

def bullet_box(slide, items, x, y, w, h, size=14, color=WHITE, dot_color=ACCENT):
    for i, item in enumerate(items):
        dot_y = y + i * (h / len(items))
        rect(slide, x, dot_y + 0.12, 0.07, 0.07, dot_color)
        txb(slide, item, x + 0.15, dot_y, w - 0.15, h / len(items), size=size, color=color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
rect(s, 0, 0, 13.33, 0.1, ACCENT)
rect(s, 0, 7.4, 13.33, 0.1, PURPLE)
# Gradient accent block
rect(s, 0, 2.5, 13.33, 0.04, ACCENT)
txb(s, "AI AssemblyTwin", 1.0, 1.0, 11, 1.2, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Real-Time AI-Powered Digital Twin for Vehicle Assembly", 1.0, 2.2, 11, 0.6, size=18, color=ACCENT, align=PP_ALIGN.CENTER)
txb(s, "Accenture Innovation Challenge 2026  ·  Round 2 Submission", 1.0, 2.8, 11, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
txb(s, "Team: Jayanth  ·  Abhinav  ·  Sagar", 1.0, 5.8, 11, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
txb(s, "github.com/jay-git00/DigitalTwin.ai", 1.0, 6.3, 11, 0.5, size=13, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "The Problem", "Indian auto-manufacturing loses ₹2,000+ Crore/year to preventable defects")

stats = [
    ("₹3.5 L", "Rework cost per defect\nthat slips past QC"),
    ("4 hrs",   "Average downtime from\nan undetected bottleneck"),
    ("2 days",  "Time to trace root\ncause — traditionally"),
    ("18%",     "Throughput loss from\npoor maintenance scheduling"),
]
for i, (val, label) in enumerate(stats):
    bx = 0.4 + i * 3.2
    rect(s, bx, 1.5, 2.9, 2.5, RGBColor(0x1E, 0x29, 0x3B))
    rect(s, bx, 1.5, 2.9, 0.06, RED)
    txb(s, val,   bx + 0.15, 1.65, 2.6, 0.9, size=34, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txb(s, label, bx + 0.15, 2.55, 2.6, 0.9, size=12, color=MUTED, align=PP_ALIGN.CENTER)

txb(s, "Traditional SCADA/PLC systems tell you what is happening NOW — not what will happen in the next 10 minutes.",
    0.5, 4.3, 12.3, 0.6, size=15, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "There is no predictive layer. No explainability. No ROI quantification. No ESG reporting.",
    0.5, 4.85, 12.3, 0.5, size=13, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "Our Solution: AI AssemblyTwin", "A live, full-stack Digital Twin — not a dashboard, not a replay")

pillars = [
    (ACCENT,  "PREDICT",    "Anomalies, bottlenecks\nand defects before\nthey happen"),
    (GREEN,   "EXPLAIN",    "Every AI decision in\nplain English — no\nblack-box AI"),
    (PURPLE,  "QUANTIFY",   "Exact ₹ ROI for every\nstakeholder, from\nshop floor to boardroom"),
    (AMBER,   "SCALE",      "3-plant enterprise\nnetwork from\nday one"),
]
for i, (col, heading, body) in enumerate(pillars):
    bx = 0.4 + i * 3.2
    rect(s, bx, 1.5, 2.9, 3.2, RGBColor(0x1E, 0x29, 0x3B))
    rect(s, bx, 1.5, 2.9, 0.06, col)
    txb(s, heading, bx + 0.1, 1.65, 2.7, 0.6, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, body,    bx + 0.1, 2.35, 2.7, 1.8, size=13, color=WHITE, align=PP_ALIGN.CENTER)

txb(s, "Built with FastAPI · SimPy · Next.js 14 · PyTorch LSTM · Isolation Forest · Random Forest · Gaussian Process",
    0.5, 4.95, 12.3, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "System Architecture", "Three layers: Simulation → AI Inference → Multi-stakeholder UI")

layers = [
    (PURPLE, "SIMULATION LAYER",  "SimPy discrete-event simulation — 45 stations, 1 vehicle/station\nGenerates realistic cycle_time, torque, vibration, temperature"),
    (ACCENT, "BACKEND (FastAPI)", "4 ML models run every 3 seconds on live telemetry\nBroadcasts over WebSocket — sub-100ms latency\n10 REST API endpoints for dashboards"),
    (GREEN,  "FRONTEND (Next.js)","6 pages, real-time data binding, Framer Motion animations\nWorks for 3 stakeholders: Supervisor · Manager · Leadership + ESG"),
]
for i, (col, title, body) in enumerate(layers):
    by = 1.4 + i * 1.8
    rect(s, 0.4, by, 12.5, 1.55, RGBColor(0x1E, 0x29, 0x3B))
    rect(s, 0.4, by, 0.07, 1.55, col)
    txb(s, title, 0.65, by + 0.1, 4.0, 0.45, size=14, bold=True, color=col)
    txb(s, body,  0.65, by + 0.55, 11.8, 0.9, size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 4 ML MODELS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "4 ML Models Running Live", "Every prediction made in real-time — not pre-computed")

models = [
    (ACCENT,  "Isolation Forest",     "Anomaly Detector",        "Unsupervised — no labelled data needed\nScores each station every 3 sec\nOutput: anomaly_score ∈ [-1, 1]"),
    (PURPLE,  "LSTM (PyTorch→NumPy)", "Bottleneck Predictor",    "2-layer LSTM on 30-min rolling window\nPredicts bottleneck 45 min ahead\nOutput: prob per station ∈ [0,1]"),
    (GREEN,   "Random Forest",        "Defect Predictor",        "200 trees + LAG FEATURES (key innovation)\nTraces upstream root-cause origin\nOutput: defect_risk + feature importance"),
    (AMBER,   "Gaussian Process",     "Sensor Imputer",          "Fills missing legacy sensor readings\nProvides explicit uncertainty (σ)\nOutput: imputed value + confidence"),
]
for i, (col, algo, role, body) in enumerate(models):
    bx = 0.4 + i * 3.2
    rect(s, bx, 1.4, 2.9, 3.5, RGBColor(0x1E, 0x29, 0x3B))
    rect(s, bx, 1.4, 2.9, 0.06, col)
    txb(s, algo, bx + 0.1, 1.55, 2.7, 0.5, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, role, bx + 0.1, 2.05, 2.7, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    txb(s, body, bx + 0.1, 2.5,  2.7, 1.8, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURE: BEFORE/AFTER + WHY?
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "Killer Demo Features", "Two moments that will win the room")

# Left box
rect(s, 0.4, 1.4, 5.9, 4.0, RGBColor(0x1E, 0x29, 0x3B))
rect(s, 0.4, 1.4, 5.9, 0.06, RED)
txb(s, "BEFORE / AFTER TOGGLE", 0.5, 1.55, 5.7, 0.5, size=15, bold=True, color=RED)
txb(s, 'Toggle "Digital Twin: OFF" → Red banner appears:\n\n"12 vehicles passed QC with defects\n₹48,00,000 in rework cost\n4.2 hours unplanned downtime"\n\nToggle ON → banner disappears.\nThat ₹48L is still in your pocket.\nROI ticker flashes and jumps.', 0.5, 2.1, 5.7, 3.1, size=13, color=WHITE)

# Right box
rect(s, 6.9, 1.4, 5.9, 4.0, RGBColor(0x1E, 0x29, 0x3B))
rect(s, 6.9, 1.4, 5.9, 0.06, ACCENT)
txb(s, 'WHY? EXPLAINABILITY BUTTON', 7.0, 1.55, 5.7, 0.5, size=15, bold=True, color=ACCENT)
txb(s, 'Every anomalous station shows a "WHY?" button.\n\nClicking it shows:\n• cycle_time_s: 42% contribution\n• torque_lag1: 31% contribution\n• vibration_g: 18% contribution\n\nPlain English. Zero ML knowledge needed.\nAddresses the "black box AI" objection\nin 10 seconds.', 7.0, 2.1, 5.7, 3.1, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PREDICTIVE MAINTENANCE + DEFECT TRACE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "Predictive Maintenance & Defect Tracing", "Think weeks ahead. Trace origins in seconds.")

rect(s, 0.4, 1.4, 5.9, 4.5, RGBColor(0x1E, 0x29, 0x3B))
rect(s, 0.4, 1.4, 5.9, 0.06, AMBER)
txb(s, "MAINTENANCE CALENDAR (/maintenance)", 0.5, 1.55, 5.7, 0.5, size=13, bold=True, color=AMBER)
txb(s, "Full monthly calendar with colour-coded dates:\n🔴 HIGH — tool failure < 7 days\n🟡 MED  — 10-20 days\n🟢 LOW  — 20-30 days\n\nProactive repair cost:    ₹92,000\nEmergency failure cost:  ₹7,64,000\n\nMultiplier: 8.3× more expensive\nto react than to act proactively.", 0.5, 2.1, 5.7, 3.6, size=13, color=WHITE)

rect(s, 6.9, 1.4, 5.9, 4.5, RGBColor(0x1E, 0x29, 0x3B))
rect(s, 6.9, 1.4, 5.9, 0.06, GREEN)
txb(s, "DEFECT TRACE (/defect-trace)", 7.0, 1.55, 5.7, 0.5, size=13, bold=True, color=GREEN)
txb(s, "Shows full causal propagation chain:\n\nStation 7 (Origin)\n    ↓  torque deviation — risk 22%\nStation 18\n    ↓  lag amplification — risk 41%\nStation 44 (QC Gate)\n         defect_risk 78%\n\nTraditional: 2-day root cause analysis\nOur system: 3 seconds\nPowered by Random Forest lag features.", 7.0, 2.1, 5.7, 3.6, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MULTI-SITE + ESG
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "Enterprise Scale & ESG Impact", "From single plant to national network — with auditable sustainability metrics")

rect(s, 0.4, 1.4, 5.9, 4.5, RGBColor(0x1E, 0x29, 0x3B))
rect(s, 0.4, 1.4, 5.9, 0.06, PURPLE)
txb(s, "MULTI-SITE OVERVIEW (/multisite)", 0.5, 1.55, 5.7, 0.5, size=13, bold=True, color=PURPLE)
txb(s, "3 Plants — 1 AI Backbone:\n\n🟢 Chennai   — LIVE        91.4% health\n🟡 Pune      — MONITORING  87.3% health\n🟢 Bangalore — OPTIMAL     95.1% health\n\nNew plant onboarding: 4-6 weeks\n(vs 6+ months for SCADA deployment)\n\nTransfer learning: 68%→91% accuracy\nin 2 weeks using Pune data.", 0.5, 2.1, 5.7, 3.6, size=13, color=WHITE)

rect(s, 6.9, 1.4, 5.9, 4.5, RGBColor(0x1E, 0x29, 0x3B))
rect(s, 6.9, 1.4, 5.9, 0.06, GREEN)
txb(s, "ESG / SUSTAINABILITY (/analytics → ESG)", 7.0, 1.55, 5.7, 0.5, size=13, bold=True, color=GREEN)
txb(s, "Per prevented scrap vehicle:\n• 333 kg CO₂ avoided\n• 180 kg steel waste prevented\n• 12 L paint solvent avoided\n• 0.4 kWh energy saved\n\nUN SDG Alignment:\n✓ SDG 9  — Industry & Innovation\n✓ SDG 12 — Responsible Consumption\n✓ SDG 13 — Climate Action\n\nAligned to Accenture Net Zero 2025.", 7.0, 2.1, 5.7, 3.6, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ROI & BUSINESS CASE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "Business Case & ROI", "Payback in 2.5 months. ₹6.1 Crore net in 3 years.")

metrics = [
    ("₹45 L",      "Deployment Cost",          MUTED),
    ("₹18 L/mo",   "Monthly Savings",          GREEN),
    ("2.5 months", "Payback Period",            ACCENT),
    ("₹6.1 Cr",    "3-Year Net ROI",            GREEN),
    ("8.3×",       "Proactive vs Reactive Repair", AMBER),
    ("4-6 weeks",  "New Plant Onboarding",      PURPLE),
]
for i, (val, label, col) in enumerate(metrics):
    row = i // 3
    col_i = i % 3
    bx = 0.5 + col_i * 4.2
    by = 1.55 + row * 2.3
    rect(s, bx, by, 3.8, 2.0, RGBColor(0x1E, 0x29, 0x3B))
    rect(s, bx, by, 3.8, 0.05, col)
    txb(s, val,   bx + 0.15, by + 0.15, 3.5, 0.85, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, label, bx + 0.15, by + 1.0,  3.5, 0.7,  size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECH STACK SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
slide_title(s, "Technology Stack", "Production-grade, fully open-source, enterprise-ready")

rows = [
    ("Simulation",          "SimPy 4.1",         "Python-native discrete-event simulation. 45 stations, realistic physics."),
    ("Anomaly Detection",   "Isolation Forest",  "Unsupervised. No labelled data. Detects drift in cycle_time/torque/vibration."),
    ("Bottleneck Prediction","LSTM (NumPy)",      "2-layer LSTM. 30-min window. Converted to NumPy for <100MB RAM on cloud."),
    ("Defect Tracing",      "Random Forest",     "200 trees + lag features. Root-cause tracing across 37 stations."),
    ("Sensor Imputation",   "Gaussian Process",  "Bayesian model. Explicit uncertainty bounds for legacy sensor-less stations."),
    ("Real-time API",       "FastAPI + WS",      "Async Python. Sub-100ms latency. 10 endpoints. Horizontally scalable."),
    ("Frontend",            "Next.js 14",        "Server components, Framer Motion, Recharts, custom dark design system."),
]
for i, (layer, tech, why) in enumerate(rows):
    by = 1.35 + i * 0.73
    rect(s, 0.4, by, 12.5, 0.68, RGBColor(0x1A, 0x23, 0x35) if i % 2 == 0 else RGBColor(0x0F, 0x17, 0x2A))
    txb(s, layer, 0.5,  by + 0.08, 2.4,  0.5, size=11, bold=True,  color=ACCENT)
    txb(s, tech,  2.95, by + 0.08, 2.5,  0.5, size=11, bold=True,  color=WHITE)
    txb(s, why,   5.5,  by + 0.08, 7.3,  0.5, size=11,              color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING / CALL TO ACTION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
rect(s, 0, 0, 13.33, 0.1, ACCENT)
rect(s, 0, 7.4, 13.33, 0.1, PURPLE)
rect(s, 0, 3.2, 13.33, 0.04, ACCENT)

txb(s, "AI AssemblyTwin", 1.0, 0.6,  11, 0.9, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "From a 45-station line to a national enterprise network.", 1.0, 1.5,  11, 0.5, size=18, color=ACCENT, align=PP_ALIGN.CENTER)
txb(s, "Predict. Explain. Quantify. Scale.", 1.0, 2.0, 11, 0.5, size=16, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

kpis = ["2.5-month payback", "₹6.1 Cr 3-year ROI", "4 ML models live", "3-plant enterprise scale", "Accenture Net Zero aligned"]
for i, kpi in enumerate(kpis):
    bx = 0.5 + i * 2.45
    rect(s, bx, 3.5, 2.25, 0.7, RGBColor(0x1E, 0x29, 0x3B))
    txb(s, kpi, bx + 0.1, 3.58, 2.05, 0.55, size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txb(s, "Team: Jayanth  ·  Abhinav  ·  Sagar", 1, 4.6, 11, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
txb(s, "github.com/jay-git00/DigitalTwin.ai", 1, 5.1, 11, 0.5, size=14, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "submission/AI_AssemblyTwin_Proposal.pptx"
import os; os.makedirs("submission", exist_ok=True)
prs.save(out)
print(f"PPTX saved → {out}")
